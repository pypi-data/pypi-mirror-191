import os
import pathlib

import click
from pptx import Presentation
from pptx.util import Inches

import logging

logger = logging.getLogger(__name__)

IMG_EXTENSIONS = [
    '.png', '.jpg', '.jpeg'
]


def collect_images(img_path):
    images = []
    for img in sorted(os.listdir(img_path)):
        if os.path.splitext(img)[1] in IMG_EXTENSIONS:
            images.append(os.path.join(img_path, img))
    return images


@click.command()
@click.option('--width', '-w', default=53.333, help='Ширина слайда в дюймах, по умолчанию: 53.333')
@click.option('--height', '-h', default=30.0, help='Высота слайда в дюймах, по умолчанию: 30')
@click.argument('img_path')
def create_pptx(img_path, width, height):
    path = pathlib.Path(img_path).absolute()
    pptx_path = path.parent / f'{path.name}.pptx'
    images = collect_images(img_path)

    if pptx_path.exists():
        prs = Presentation(pptx_path)

        for i, img in enumerate(images):
            slide = prs.slides[i]

            # Remove old image
            pic = slide.shapes[0]
            pic = pic._element
            pic.getparent().remove(pic)

            # Add new image
            pic = slide.shapes.add_picture(
                img, 0, 0,
                width=prs.slide_width,
                height=prs.slide_height
            )
    else:
        prs = Presentation()
        prs.slide_width = Inches(width)
        prs.slide_height = Inches(height)
        blank_slide_layout = prs.slide_layouts[6]

        for img in images:
            slide = prs.slides.add_slide(blank_slide_layout)

            pic = slide.shapes.add_picture(
                img, 0, 0,
                width=prs.slide_width,
                height=prs.slide_height
            )

            # notes_slide = slide.notes_slide
            # text_frame = notes_slide.notes_text_frame
            # text_frame.text = ''  # место для заметок докладчика

    if len(images) > 0:
        prs.save(pptx_path)
    else:
        logger.error("Empty images set. File not created.")


if __name__ == '__main__':
    create_pptx()
